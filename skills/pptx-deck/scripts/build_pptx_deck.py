#!/usr/bin/env python3
"""Build a simple PPTX deck (title, data table, chart) from a CSV.

Usage:
    python3 build_pptx_deck.py path/to/data.csv --title "My Deck" --out deck.pptx
"""
import argparse
import sys
import tempfile
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

MAX_TABLE_ROWS = 10


def pick_numeric_column(df: pd.DataFrame):
    numeric = df.select_dtypes(include="number").columns.tolist()
    return numeric[0] if numeric else None


def pick_category_column(df: pd.DataFrame, exclude: str, max_categories: int = 12):
    for col in df.columns:
        if col == exclude:
            continue
        if df[col].dtype == object and 1 < df[col].nunique(dropna=True) <= max_categories:
            return col
    return None


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_table_slide(prs: Presentation, df: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Data Preview (showing {min(len(df), MAX_TABLE_ROWS)} of {len(df)} rows)"

    preview = df.head(MAX_TABLE_ROWS)
    rows, cols = preview.shape[0] + 1, preview.shape[1]
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for c, col_name in enumerate(preview.columns):
        table.cell(0, c).text = str(col_name)

    for r in range(preview.shape[0]):
        for c in range(preview.shape[1]):
            table.cell(r + 1, c).text = str(preview.iat[r, c])

    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)


def add_chart_slide(prs: Presentation, df: pd.DataFrame, numeric_col: str, category_col):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = (
        f"{numeric_col} by {category_col}" if category_col else f"Distribution of {numeric_col}"
    )

    fig, ax = plt.subplots(figsize=(8, 4.5))
    if category_col:
        grouped = df.groupby(category_col)[numeric_col].sum().sort_values(ascending=False)
        grouped.plot(kind="bar", ax=ax, color="#2E86AB")
    else:
        df[numeric_col].plot(kind="hist", ax=ax, color="#2E86AB", bins=20)
    fig.tight_layout()

    with tempfile.TemporaryDirectory() as tmp:
        chart_path = Path(tmp) / "chart.png"
        fig.savefig(chart_path, dpi=150)
        plt.close(fig)
        slide.shapes.add_picture(str(chart_path), Inches(0.75), Inches(1.5), width=Inches(8.5))


def build_deck(csv_path: str, title: str, out_path: str) -> None:
    df = pd.read_csv(csv_path)
    prs = Presentation()

    add_title_slide(prs, title, f"Generated from {Path(csv_path).name} — {len(df)} rows, {len(df.columns)} columns")
    add_table_slide(prs, df)

    numeric_col = pick_numeric_column(df)
    if numeric_col:
        category_col = pick_category_column(df, exclude=numeric_col)
        add_chart_slide(prs, df, numeric_col, category_col)

    prs.save(out_path)
    print(f"Wrote {out_path}")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("csv_path")
    parser.add_argument("--title", default="Data Deck")
    parser.add_argument("--out", default="deck.pptx")
    args = parser.parse_args()

    try:
        build_deck(args.csv_path, args.title, args.out)
    except FileNotFoundError:
        print(f"File not found: {args.csv_path}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
